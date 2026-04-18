#!/usr/bin/env python3
"""Convert a PDF slide deck to PPTX with each page as a full-slide image.

Usage:
    uv run python scripts/pdf_to_pptx.py input.pdf [output.pptx]
"""

import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def pdf_to_pptx(pdf_path: str, pptx_path: str | None = None):
    pdf = Path(pdf_path)
    if pptx_path is None:
        pptx_path = str(pdf.with_suffix(".pptx"))

    # Convert PDF pages to PNG using pdftoppm (poppler)
    with tempfile.TemporaryDirectory() as tmpdir:
        print(f"Converting {pdf} to images...")
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", "300", str(pdf), f"{tmpdir}/page"],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            print(f"pdftoppm failed: {result.stderr}")
            sys.exit(1)

        pages = sorted(Path(tmpdir).glob("page-*.png"))
        print(f"  {len(pages)} pages extracted")

        # Create PPTX (16:9 widescreen)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank layout

        for i, img_path in enumerate(pages):
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(img_path),
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )
            print(f"  Page {i+1}/{len(pages)}")

        prs.save(pptx_path)
        print(f"Saved: {pptx_path}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: uv run python scripts/pdf_to_pptx.py input.pdf [output.pptx]")
        sys.exit(1)
    pdf_to_pptx(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
